import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not installed yet. Please wait.")
    sys.exit(1)

def build_presentation():
    prs = Presentation()
    # Set standard widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Theme colors
    COLOR_BG = RGBColor(0xE5, 0xF2, 0xC9)       # Cream
    COLOR_CARD = RGBColor(0xFF, 0xFF, 0xFF)     # White
    COLOR_PRIMARY = RGBColor(0x7F, 0x53, 0x4B)  # Terracotta
    COLOR_SECONDARY = RGBColor(0x8C, 0x70, 0x5F)# Brown
    COLOR_HEADING = RGBColor(0x1F, 0x03, 0x18)  # Near black
    COLOR_TEXT = RGBColor(0x1E, 0x1A, 0x1D)     # Text

    blank_layout = prs.slide_layouts[6] # Blank slide

    def set_background(slide):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = COLOR_BG
        rect.line.fill.background()

    def add_card_base(slide):
        left = Inches(0.4)
        top = Inches(0.4)
        width = Inches(12.53)
        height = Inches(6.7)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1)

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
        accent.fill.solid()
        accent.fill.fore_color.rgb = COLOR_PRIMARY
        accent.line.fill.background()

    def add_header(slide, category, title):
        # Category
        cat_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(10), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = 'Trebuchet MS'
        cat_p.font.size = Pt(10)
        cat_p.font.bold = True
        cat_p.font.color.rgb = COLOR_PRIMARY

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.9), Inches(11.5), Inches(0.6))
        title_tf = title_box.text_frame
        title_p = title_tf.paragraphs[0]
        title_p.text = title
        title_p.font.name = 'Trebuchet MS'
        title_p.font.size = Pt(22)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_HEADING

    def add_footer(slide):
        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.6), Inches(10), Inches(0.3))
        footer_tf = footer_box.text_frame
        footer_p = footer_tf.paragraphs[0]
        footer_p.text = "MathSecure - Sistem Autentikasi Kontekstual | Politeknik Caltex Riau"
        footer_p.font.name = 'Arial'
        footer_p.font.size = Pt(8)
        footer_p.font.color.rgb = COLOR_SECONDARY

    def make_content_slide(category, title):
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide)
        add_card_base(slide)
        add_header(slide, category, title)
        add_footer(slide)
        return slide

    # ==========================================================
    # SLIDE 1: Judul
    # ==========================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1)
    # Background card frame
    card1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.73), Inches(5.9))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_SECONDARY
    card1.line.width = Pt(2)

    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.73), Inches(0.12))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = COLOR_PRIMARY
    accent1.line.fill.background()

    # Category Tag
    tag_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(10.9), Inches(0.4))
    tag_p = tag_box.text_frame.paragraphs[0]
    tag_p.text = "STUDI KASUS MATEMATIKA DISKRIT"
    tag_p.font.name = 'Trebuchet MS'
    tag_p.font.size = Pt(12)
    tag_p.font.bold = True
    tag_p.font.color.rgb = COLOR_PRIMARY

    # Main Title
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.6))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Sistem Autentikasi Berbasis React\ndan Kriptografi XOR"
    title_p.font.name = 'Trebuchet MS'
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_HEADING

    # Subtitle
    sub_box = slide1.shapes.add_textbox(Inches(1.2), Inches(3.6), Inches(10.9), Inches(0.6))
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = "Media Pembelajaran Logika Proposisional, Teori Himpunan, Relasi Biner & Hashing SHA-256"
    sub_p.font.name = 'Arial'
    sub_p.font.size = Pt(14)
    sub_p.font.color.rgb = COLOR_SECONDARY
    sub_p.font.italic = True

    # Members list label
    m_lbl_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(10.9), Inches(0.3))
    m_lbl_p = m_lbl_box.text_frame.paragraphs[0]
    m_lbl_p.text = "Oleh Kelompok:"
    m_lbl_p.font.name = 'Arial'
    m_lbl_p.font.size = Pt(11)
    m_lbl_p.font.bold = True
    m_lbl_p.font.color.rgb = COLOR_SECONDARY

    # Members names grid-style
    m_names_box = slide1.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(10.9), Inches(0.5))
    m_names_p = m_names_box.text_frame.paragraphs[0]
    m_names_p.text = "1. M Risqullah Naufal    |    2. Lingzhi    |    3. Rafli Hanafi    |    4. Victor"
    m_names_p.font.name = 'Arial'
    m_names_p.font.size = Pt(13)
    m_names_p.font.bold = True
    m_names_p.font.color.rgb = COLOR_TEXT

    # Footer
    univ_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.4))
    univ_p = univ_box.text_frame.paragraphs[0]
    univ_p.text = "Program Studi Teknik Informatika - Politeknik Caltex Riau  |  TA 2025/2026"
    univ_p.font.name = 'Arial'
    univ_p.font.size = Pt(11)
    univ_p.font.color.rgb = COLOR_SECONDARY

    # ==========================================================
    # SLIDE 2: Latar Belakang & Masalah
    # ==========================================================
    slide2 = make_content_slide("BAB I - PENDAHULUAN", "Latar Belakang & Rumusan Masalah")

    # Column 1 text box
    col1_box = slide2.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(5.5), Inches(4.5))
    col1_tf = col1_box.text_frame
    col1_tf.word_wrap = True
    
    p = col1_tf.paragraphs[0]
    p.text = "Latar Belakang Proyek"
    p.font.name = 'Trebuchet MS'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.space_after = Pt(14)

    bullets = [
        "Pentingnya aspek keamanan siber pada sistem otentikasi akun web modern.",
        "Konsep dasar pendukung keamanan (Logika Proposisional, Himpunan, Relasi) sering diajarkan secara sangat teoritis di perkuliahan.",
        "Kebutuhan media simulasi interaktif agar mahasiswa dapat memanipulasi parameter boolean/matriks secara praktis dan visual."
    ]
    for b in bullets:
        bp = col1_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # Column 2 container
    col2_rect = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.5))
    col2_rect.fill.solid()
    col2_rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
    col2_rect.line.color.rgb = RGBColor(0xE2, 0xD9, 0xD2)

    col2_box = slide2.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.2), Inches(4.1))
    col2_tf = col2_box.text_frame
    col2_tf.word_wrap = True
    
    p2 = col2_tf.paragraphs[0]
    p2.text = "Rumusan Masalah Utama"
    p2.font.name = 'Trebuchet MS'
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_HEADING
    p2.space_after = Pt(14)

    q_items = [
        "1. Bagaimana memodelkan kelayakan login berdasarkan parameter lingkungan menggunakan Aljabar Boolean?",
        "2. Bagaimana memetakan hak otorisasi pengguna (Admin & Dosen) secara visual melalui Diagram Venn?",
        "3. Bagaimana memetakan relasi peran pengguna dalam bentuk representasi biner matriks?",
        "4. Bagaimana memperlihatkan tahapan enkripsi bitwise XOR dan hash SHA-256?"
    ]
    for q in q_items:
        qp = col2_tf.add_paragraph()
        qp.text = q
        qp.font.name = 'Arial'
        qp.font.size = Pt(11)
        qp.font.color.rgb = COLOR_TEXT
        qp.space_after = Pt(12)

    # ==========================================================
    # SLIDE 3: Tujuan & Manfaat
    # ==========================================================
    slide3 = make_content_slide("BAB I - PENDAHULUAN", "Tujuan & Manfaat Proyek")

    # Box 1: Tujuan
    t_rect = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.6), Inches(5.6), Inches(4.5))
    t_rect.fill.solid()
    t_rect.fill.fore_color.rgb = COLOR_CARD
    t_rect.line.color.rgb = COLOR_PRIMARY
    t_rect.line.width = Pt(2)

    t_box = slide3.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.1))
    t_tf = t_box.text_frame
    t_tf.word_wrap = True
    tp = t_tf.paragraphs[0]
    tp.text = "Tujuan Penelitian"
    tp.font.name = 'Trebuchet MS'
    tp.font.size = Pt(16)
    tp.font.bold = True
    tp.font.color.rgb = COLOR_PRIMARY
    tp.space_after = Pt(14)

    t_bullets = [
        "Menerapkan ekspresi logika proposisional pada evaluasi login web.",
        "Membuat modul pembelajaran interaktif berupa visualisasi Diagram Venn untuk materi Teori Himpunan.",
        "Membangun representasi matriks relasi user-role biner.",
        "Menguji modul simulator bitwise XOR dan enkripsi SHA-256."
    ]
    for b in t_bullets:
        bp = t_tf.add_paragraph()
        bp.text = "✓  " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # Box 2: Manfaat
    m_rect = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.5))
    m_rect.fill.solid()
    m_rect.fill.fore_color.rgb = COLOR_CARD
    m_rect.line.color.rgb = COLOR_SECONDARY
    m_rect.line.width = Pt(2)

    m_box = slide3.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.1))
    m_tf = m_box.text_frame
    m_tf.word_wrap = True
    mp = m_tf.paragraphs[0]
    mp.text = "Manfaat Yang Dihasilkan"
    mp.font.name = 'Trebuchet MS'
    mp.font.size = Pt(16)
    mp.font.bold = True
    mp.font.color.rgb = COLOR_SECONDARY
    mp.space_after = Pt(14)

    m_bullets = [
        "Bagi Mahasiswa/Akademik:\nMenjembatani jurang pemisah antara teori matematika diskrit yang abstrak dan implementasi perangkat lunak nyata.",
        "Bagi Pengembang/Keamanan:\nMenunjukkan efektivitas Context-Aware Security yang mengurangi gesekan login bagi pengguna terpercaya tanpa menurunkan level keamanan."
    ]
    for b in m_bullets:
        bp = m_tf.add_paragraph()
        bp.text = "✦  " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(12)

    # ==========================================================
    # SLIDE 4: Logika Proposisional
    # ==========================================================
    slide4 = make_content_slide("BAB II & IV - LOGIKA MATEMATIKA", "Logika Proposisional Otentikasi")

    intro_box = slide4.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(0.4))
    intro_p = intro_box.text_frame.paragraphs[0]
    intro_p.text = "Keputusan keberhasilan login (L) dievaluasi secara dinamis berdasarkan 4 variabel proposisional:"
    intro_p.font.name = 'Arial'
    intro_p.font.size = Pt(12)
    intro_p.font.color.rgb = COLOR_TEXT

    # 4 Cards representing variables
    vars_data = [
        {"tag": "p (Kredensial)", "desc": "Username, password, dan pilihan role sesuai database."},
        {"tag": "q (Akun Aktif)", "desc": "Status akun pengguna tidak ditangguhkan."},
        {"tag": "r (Perangkat)", "desc": "Perangkat terdaftar ('Remember Device' aktif)."},
        {"tag": "s (Captcha)", "desc": "Pengguna berhasil menjawab soal Captcha."}
    ]

    for idx, v in enumerate(vars_data):
        xPos = Inches(0.7 + (idx * 2.9))
        rect = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, xPos, Inches(2.0), Inches(2.7), Inches(2.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
        rect.line.color.rgb = RGBColor(0xE2, 0xD9, 0xD2)

        v_box = slide4.shapes.add_textbox(xPos + Inches(0.1), Inches(2.1), Inches(2.5), Inches(1.8))
        v_tf = v_box.text_frame
        v_tf.word_wrap = True
        vp = v_tf.paragraphs[0]
        vp.text = v["tag"]
        vp.font.name = 'Trebuchet MS'
        vp.font.size = Pt(14)
        vp.font.bold = True
        vp.font.color.rgb = COLOR_PRIMARY
        vp.space_after = Pt(10)

        vdesc = v_tf.add_paragraph()
        vdesc.text = v["desc"]
        vdesc.font.name = 'Arial'
        vdesc.font.size = Pt(10)
        vdesc.font.color.rgb = COLOR_TEXT

    # Formula Box
    f_rect = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.3), Inches(11.5), Inches(1.8))
    f_rect.fill.solid()
    f_rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
    f_rect.line.color.rgb = COLOR_PRIMARY
    f_rect.line.width = Pt(1.5)

    f_box = slide4.shapes.add_textbox(Inches(0.9), Inches(4.4), Inches(11.1), Inches(1.6))
    f_tf = f_box.text_frame
    f_tf.word_wrap = True
    
    fp1 = f_tf.paragraphs[0]
    fp1.text = "FORMULA ALJABAR BOOLEAN (L = STATUS HASIL LOGIN)"
    fp1.font.name = 'Trebuchet MS'
    fp1.font.size = Pt(11)
    fp1.font.bold = True
    fp1.font.color.rgb = COLOR_SECONDARY
    fp1.space_after = Pt(6)

    fp2 = f_tf.add_paragraph()
    fp2.text = "L = (p ∧ q ∧ r) ∨ (p ∧ q ∧ ¬r ∧ s)"
    fp2.font.name = 'Courier New'
    fp2.font.size = Pt(24)
    fp2.font.bold = True
    fp2.font.color.rgb = COLOR_HEADING
    fp2.space_after = Pt(6)

    fp3 = f_tf.add_paragraph()
    fp3.text = "*Keterangan: Login berhasil jika kredensial valid (p) dan akun aktif (q), DAN (perangkat dikenali (r) ATAU jawaban captcha benar (s))."
    fp3.font.name = 'Arial'
    fp3.font.size = Pt(9)
    fp3.font.color.rgb = COLOR_TEXT
    fp3.font.italic = True

    # ==========================================================
    # SLIDE 5: Tabel Kebenaran
    # ==========================================================
    slide5 = make_content_slide("BAB IV - PEMBAHASAN", "Tabel Kebenaran Otentikasi")

    tbl_intro = slide5.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4))
    tbl_intro.text_frame.paragraphs[0].text = "Representasi tabel kebenaran yang diimplementasikan pada mesin logika proposisional:"
    tbl_intro.text_frame.paragraphs[0].font.name = 'Arial'
    tbl_intro.text_frame.paragraphs[0].font.size = Pt(11)
    tbl_intro.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    table_data = [
        ["p", "q", "r", "s", "p ∧ q ∧ r", "p ∧ q ∧ ¬r ∧ s", "Hasil (L)", "Kasus Keterangan"],
        ["True", "True", "True", "Apapun", "True", "False", "True", "SUKSES: Login di Perangkat Terpercaya"],
        ["True", "True", "False", "True", "False", "True", "True", "SUKSES: Login lewat Verifikasi Captcha"],
        ["True", "True", "False", "False", "False", "False", "False", "GAGAL: Captcha tidak cocok/valid"],
        ["True", "False", "Apapun", "Apapun", "False", "False", "False", "GAGAL: Akun ditangguhkan (Suspended)"],
        ["False", "Apapun", "Apapun", "Apapun", "False", "False", "False", "GAGAL: Kredensial tidak cocok"]
    ]

    table_shape = slide5.shapes.add_table(6, 8, Inches(0.7), Inches(1.9), Inches(11.5), Inches(4.0))
    table = table_shape.table

    # Format column widths
    widths = [Inches(0.8), Inches(0.8), Inches(0.8), Inches(1.0), Inches(1.5), Inches(1.8), Inches(1.2), Inches(3.6)]
    for idx, w in enumerate(widths):
        table.columns[idx].width = w

    for r_idx, row in enumerate(table_data):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = COLOR_HEADING
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
            else:
                p.font.color.rgb = COLOR_TEXT
                if c_idx == 6: # Hasil (L) column
                    p.font.bold = True
                    if val == "True":
                        p.font.color.rgb = RGBColor(0x04, 0x78, 0x57) # emerald green
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
                    else:
                        p.font.color.rgb = RGBColor(0xBE, 0x12, 0x3C) # rose red
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF1, 0xF2)

    # ==========================================================
    # SLIDE 6: Decision Tree
    # ==========================================================
    slide6 = make_content_slide("BAB IV - PEMBAHASAN", "Pohon Keputusan (Decision Tree) Login")

    # Column 1
    t_box2 = slide6.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(4.5))
    t2_tf = t_box2.text_frame
    t2_tf.word_wrap = True
    tp2 = t2_tf.paragraphs[0]
    tp2.text = "Alur Evaluasi Pohon Keputusan"
    tp2.font.name = 'Trebuchet MS'
    tp2.font.size = Pt(16)
    tp2.font.bold = True
    tp2.font.color.rgb = COLOR_PRIMARY
    tp2.space_after = Pt(14)

    tree_bullets = [
        "Pohon keputusan menyederhanakan alur proposisi biner menjadi serangkaian langkah percabangan terurut.",
        "Di antarmuka web, pohon keputusan digambar menggunakan grafis interaktif.",
        "Jalur evaluasi biner akan menyala secara dinamis: node bernilai True diwarnai Hijau, dan node bernilai False diwarnai Merah.",
        "Memudahkan proses evaluasi visual dan pengajaran logika sekuensial bagi mahasiswa IT."
    ]
    for b in tree_bullets:
        bp = t2_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(12)

    # Column 2 (ASCII Tree container)
    ascii_rect = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.5))
    ascii_rect.fill.solid()
    ascii_rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
    ascii_rect.line.color.rgb = RGBColor(0xE2, 0xD9, 0xD2)

    ascii_box = slide6.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.2), Inches(4.3))
    ascii_tf = ascii_box.text_frame
    ascii_tf.word_wrap = True
    ap = ascii_tf.paragraphs[0]
    ap.text = (
        "             [ INPUT DATA LOGIN ]\n"
        "                       │\n"
        "             Kredensial Valid? (p)\n"
        "              ├───(Tidak)───► [LOGIN GAGAL]\n"
        "              └───(Ya)\n"
        "                    │\n"
        "               Akun Aktif? (q)\n"
        "                ├───(Tidak)───► [LOGIN GAGAL]\n"
        "                └───(Ya)\n"
        "                      │\n"
        "                 Ingat Perangkat? (r)\n"
        "                  ├───(Ya)─────────────► [SUKSES]\n"
        "                  └───(Tidak)\n"
        "                        │\n"
        "                   Captcha Valid? (s)\n"
        "                    ├───(Ya)───────────► [SUKSES]\n"
        "                    └───(Tidak)────────► [GAGAL]"
    )
    ap.font.name = 'Courier New'
    ap.font.size = Pt(10)
    ap.font.bold = True
    ap.font.color.rgb = COLOR_HEADING

    # ==========================================================
    # SLIDE 7: Teori Himpunan & Notasi Peran
    # ==========================================================
    slide7 = make_content_slide("BAB II & IV - TEORI HIMPUNAN", "Teori Himpunan & Peran Pengguna")

    intro_set = slide7.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4))
    intro_set.text_frame.paragraphs[0].text = "Otorisasi pengguna dipetakan sebagai anggota dari himpunan bagian dari semesta (U):"
    intro_set.text_frame.paragraphs[0].font.name = 'Arial'
    intro_set.text_frame.paragraphs[0].font.size = Pt(11)
    intro_set.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    sets_info = [
        {"name": "Himpunan Semesta (U)", "val": "U = { Budi, Siti, Dewi, Gede, Anto }", "desc": "Semua pengguna terdaftar di sistem."},
        {"name": "Himpunan Admin (A)", "val": "A = { Budi, Gede, Anto }", "desc": "Pengguna dengan peran Admin."},
        {"name": "Himpunan Dosen (D)", "val": "D = { Budi, Siti, Dewi }", "desc": "Pengguna dengan peran Dosen."}
    ]

    for idx, s in enumerate(sets_info):
        xPos = Inches(0.7 + (idx * 3.9))
        rect = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, xPos, Inches(1.9), Inches(3.7), Inches(2.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
        rect.line.color.rgb = RGBColor(0xE2, 0xD9, 0xD2)

        s_box = slide7.shapes.add_textbox(xPos + Inches(0.1), Inches(2.0), Inches(3.5), Inches(1.8))
        s_tf = s_box.text_frame
        s_tf.word_wrap = True
        sp = s_tf.paragraphs[0]
        sp.text = s["name"]
        sp.font.name = 'Trebuchet MS'
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.color.rgb = COLOR_PRIMARY
        sp.space_after = Pt(10)

        sval = s_tf.add_paragraph()
        sval.text = s["val"]
        sval.font.name = 'Courier New'
        sval.font.size = Pt(11)
        sval.font.bold = True
        sval.font.color.rgb = COLOR_HEADING
        sval.space_after = Pt(8)

        sdesc = s_tf.add_paragraph()
        sdesc.text = s["desc"]
        sdesc.font.name = 'Arial'
        sdesc.font.size = Pt(9.5)
        sdesc.font.color.rgb = COLOR_TEXT

    # Set Operations Box
    sop_rect = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.2), Inches(11.5), Inches(1.9))
    sop_rect.fill.solid()
    sop_rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
    sop_rect.line.color.rgb = COLOR_PRIMARY
    sop_rect.line.width = Pt(1.5)

    sop_box = slide7.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(11.1), Inches(1.7))
    sop_tf = sop_box.text_frame
    sop_tf.word_wrap = True
    
    sopp1 = sop_tf.paragraphs[0]
    sopp1.text = "OPERASI ALJABAR HIMPUNAN PADA GRUP OTORISASI"
    sopp1.font.name = 'Trebuchet MS'
    sopp1.font.size = Pt(11)
    sopp1.font.bold = True
    sopp1.font.color.rgb = COLOR_SECONDARY
    sopp1.space_after = Pt(8)

    sop_bullets = [
        "Irisan (A ∩ D) = { Budi }  --> Pengguna dengan peran ganda Admin sekaligus Dosen.",
        "Selisih (A - D) = { Gede, Anto }  --> Pengguna dengan peran murni Admin saja.",
        "Selisih (D - A) = { Siti, Dewi }  --> Pengguna dengan peran murni Dosen saja."
    ]
    for b in sop_bullets:
        bp = sop_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Courier New'
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = COLOR_HEADING
        bp.space_after = Pt(6)

    # ==========================================================
    # SLIDE 8: Visualisasi Diagram Venn
    # ==========================================================
    slide8 = make_content_slide("BAB IV - PEMBAHASAN", "Visualisasi Diagram Venn")

    # Column 1
    v_box = slide8.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(4.5))
    v_tf = v_box.text_frame
    v_tf.word_wrap = True
    vp = v_tf.paragraphs[0]
    vp.text = "Diagram Venn Interaktif"
    vp.font.name = 'Trebuchet MS'
    vp.font.size = Pt(16)
    vp.font.bold = True
    vp.font.color.rgb = COLOR_PRIMARY
    vp.space_after = Pt(14)

    v_bullets = [
        "Direpresentasikan dengan grafik SVG dinamis di antarmuka sandbox web.",
        "Area perpotongan lingkaran melambangkan Irisan (A ∩ D), yang memuat user Budi.",
        "Area lingkaran Admin saja memuat Gede dan Anto.",
        "Area lingkaran Dosen saja memuat Siti dan Dewi.",
        "Visualisasi mempermudah peninjauan struktur set data database relasional secara intuitif."
    ]
    for b in v_bullets:
        bp = v_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(12)

    # Column 2 Venn Circles
    xCenter = 9.2
    yCenter = 3.6
    
    # Circle 1 (Admin)
    c1 = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xCenter - 1.5), Inches(yCenter - 1.2), Inches(2.2), Inches(2.2))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_PRIMARY
    c1.fill.transparency = 0.85
    c1.line.color.rgb = COLOR_PRIMARY
    c1.line.width = Pt(2)
    
    c1_txt = slide8.shapes.add_textbox(Inches(xCenter - 1.4), Inches(yCenter - 1.8), Inches(1.5), Inches(0.4))
    c1_txt.text_frame.paragraphs[0].text = "Admin (A)"
    c1_txt.text_frame.paragraphs[0].font.name = 'Trebuchet MS'
    c1_txt.text_frame.paragraphs[0].font.bold = True
    c1_txt.text_frame.paragraphs[0].font.size = Pt(12)
    c1_txt.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY

    # Circle 2 (Dosen)
    c2 = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xCenter - 0.4), Inches(yCenter - 1.2), Inches(2.2), Inches(2.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_SECONDARY
    c2.fill.transparency = 0.85
    c2.line.color.rgb = COLOR_SECONDARY
    c2.line.width = Pt(2)

    c2_txt = slide8.shapes.add_textbox(Inches(xCenter + 0.8), Inches(yCenter - 1.8), Inches(1.5), Inches(0.4))
    c2_txt.text_frame.paragraphs[0].text = "Dosen (D)"
    c2_txt.text_frame.paragraphs[0].font.name = 'Trebuchet MS'
    c2_txt.text_frame.paragraphs[0].font.bold = True
    c2_txt.text_frame.paragraphs[0].font.size = Pt(12)
    c2_txt.text_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY

    # Add text labels inside circles
    l1_box = slide8.shapes.add_textbox(Inches(xCenter - 1.2), Inches(yCenter - 0.3), Inches(0.8), Inches(0.8))
    l1_tf = l1_box.text_frame
    l1_p = l1_tf.paragraphs[0]
    l1_p.text = "gede\nanto"
    l1_p.font.name = 'Arial'
    l1_p.font.size = Pt(10)
    l1_p.font.bold = True
    l1_p.font.color.rgb = COLOR_TEXT
    l1_p.alignment = PP_ALIGN.CENTER

    l2_box = slide8.shapes.add_textbox(Inches(xCenter - 0.2), Inches(yCenter - 0.1), Inches(0.8), Inches(0.5))
    l2_tf = l2_box.text_frame
    l2_p = l2_tf.paragraphs[0]
    l2_p.text = "budi"
    l2_p.font.name = 'Arial'
    l2_p.font.size = Pt(11)
    l2_p.font.bold = True
    l2_p.font.color.rgb = COLOR_PRIMARY
    l2_p.alignment = PP_ALIGN.CENTER

    l3_box = slide8.shapes.add_textbox(Inches(xCenter + 0.8), Inches(yCenter - 0.3), Inches(0.8), Inches(0.8))
    l3_tf = l3_box.text_frame
    l3_p = l3_tf.paragraphs[0]
    l3_p.text = "siti\ndewi"
    l3_p.font.name = 'Arial'
    l3_p.font.size = Pt(10)
    l3_p.font.bold = True
    l3_p.font.color.rgb = COLOR_TEXT
    l3_p.alignment = PP_ALIGN.CENTER

    # ==========================================================
    # SLIDE 9: Relasi & Matriks Relasi
    # ==========================================================
    slide9 = make_content_slide("BAB II & IV - RELASI MATEMATIKA", "Relasi & Matriks Relasi Biner")

    # Column 1
    r_box = slide9.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.8), Inches(4.5))
    r_tf = r_box.text_frame
    r_tf.word_wrap = True
    rp = r_tf.paragraphs[0]
    rp.text = "Relasi Keanggotaan Peran"
    rp.font.name = 'Trebuchet MS'
    rp.font.size = Pt(16)
    rp.font.bold = True
    rp.font.color.rgb = COLOR_PRIMARY
    rp.space_after = Pt(14)

    r_bullets = [
        "Relasi biner R memetakan Pengguna ke Peran (Role) otorisasi mereka.",
        "Dinyatakan dalam himpunan pasangan terurut:\nR = { (Budi, Admin), (Budi, Dosen), (Siti, Dosen), (Dewi, Dosen), (Gede, Admin), (Anto, Admin) }",
        "Tipe Relasi: One-to-Many. Relasi biner ini bukan fungsi karena elemen domain Budi dipetakan ke lebih dari satu kodomain."
    ]
    for b in r_bullets:
        bp = r_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(12)

    # Column 2
    r_rect = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.6), Inches(4.5))
    r_rect.fill.solid()
    r_rect.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
    r_rect.line.color.rgb = RGBColor(0xE2, 0xD9, 0xD2)

    m_title = slide9.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.2), Inches(0.4))
    m_title.text_frame.paragraphs[0].text = "Representasi Matriks Relasi Biner (M)"
    m_title.text_frame.paragraphs[0].font.name = 'Trebuchet MS'
    m_title.text_frame.paragraphs[0].font.bold = True
    m_title.text_frame.paragraphs[0].font.size = Pt(14)
    m_title.text_frame.paragraphs[0].font.color.rgb = COLOR_HEADING

    m_table_shape = slide9.shapes.add_table(6, 3, Inches(7.0), Inches(2.2), Inches(5.2), Inches(3.2))
    m_table = m_table_shape.table

    # Format table widths
    m_table.columns[0].width = Inches(2.2)
    m_table.columns[1].width = Inches(1.5)
    m_table.columns[2].width = Inches(1.5)

    m_data = [
        ["Pengguna", "Admin", "Dosen"],
        ["Budi", "1", "1"],
        ["Siti", "0", "1"],
        ["Dewi", "0", "1"],
        ["Gede", "1", "0"],
        ["Anto", "1", "0"]
    ]

    for r_idx, row in enumerate(m_data):
        for c_idx, val in enumerate(row):
            cell = m_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(11)
            p.alignment = PP_ALIGN.CENTER
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = COLOR_HEADING
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
            else:
                p.font.color.rgb = COLOR_TEXT
                if c_idx > 0:
                    p.font.bold = True
                    if val == "1":
                        p.font.color.rgb = COLOR_PRIMARY

    # ==========================================================
    # SLIDE 10: Kriptografi Password
    # ==========================================================
    slide10 = make_content_slide("BAB II & IV - KRIPTOGRAFI", "Kriptografi: SHA-256 vs XOR Cipher")

    # Column 1
    sha_rect = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.6), Inches(5.6), Inches(4.5))
    sha_rect.fill.solid()
    sha_rect.fill.fore_color.rgb = COLOR_CARD
    sha_rect.line.color.rgb = COLOR_PRIMARY
    sha_rect.line.width = Pt(2)

    sha_box = slide10.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.1))
    sha_tf = sha_box.text_frame
    sha_tf.word_wrap = True
    
    shp = sha_tf.paragraphs[0]
    shp.text = "1. Hashing SHA-256"
    shp.font.name = 'Trebuchet MS'
    shp.font.size = Pt(16)
    shp.font.bold = True
    shp.font.color.rgb = COLOR_PRIMARY
    shp.space_after = Pt(12)

    sha_bullets = [
        "Fungsi kriptografi satu-arah (One-way Hashing) standar industri yang aman.",
        "Mengonversi input password menjadi digest unik sepanjang 256-bit (64 karakter heksadesimal).",
        "Sifat: Non-reversible (tidak dapat didekripsi balik ke teks asli).",
        "Proses: Padding bit, konstanta inisialisasi register akar pecahan bilangan prima, kompresi 64 putaran."
    ]
    for b in sha_bullets:
        bp = sha_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # Column 2
    xor_rect = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(4.5))
    xor_rect.fill.solid()
    xor_rect.fill.fore_color.rgb = COLOR_CARD
    xor_rect.line.color.rgb = COLOR_SECONDARY
    xor_rect.line.width = Pt(2)

    xor_box = slide10.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.1))
    xor_tf = xor_box.text_frame
    xor_tf.word_wrap = True
    
    xop = xor_tf.paragraphs[0]
    xop.text = "2. Enkripsi XOR Simetris"
    xop.font.name = 'Trebuchet MS'
    xop.font.size = Pt(16)
    xop.font.bold = True
    xop.font.color.rgb = COLOR_SECONDARY
    xop.space_after = Pt(12)

    xor_bullets = [
        "Digunakan untuk visualisasi akademis bitwise biner di laboratorium interaktif.",
        "Meng-XOR bitwise setiap karakter kata sandi dengan Kunci K = 90 (01011010).",
        "Sifat: Involutif. Operasi enkripsi dan dekripsi menggunakan fungsi biner matematika yang persis sama.",
        "Formula: P ⊕ K = C  dan  C ⊕ K = P."
    ]
    for b in xor_bullets:
        bp = xor_tf.add_paragraph()
        bp.text = "• " + b
        bp.font.name = 'Arial'
        bp.font.size = Pt(11)
        bp.font.color.rgb = COLOR_TEXT
        bp.space_after = Pt(10)

    # ==========================================================
    # SLIDE 11: Skenario Pengujian
    # ==========================================================
    slide11 = make_content_slide("BAB V - HASIL & ANALISIS", "Skenario Pengujian & Hasil")

    test_intro = slide11.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(0.4))
    test_intro.text_frame.paragraphs[0].text = "Sistem diuji dengan 5 skenario dengan hasil keputusan logika proposisional yang akurat:"
    test_intro.text_frame.paragraphs[0].font.name = 'Arial'
    test_intro.text_frame.paragraphs[0].font.size = Pt(11)
    test_intro.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    test_data = [
        ["No", "Skenario Pengujian Kondisi (p, q, r, s)", "Hasil Ekspektasi", "Hasil Aktual", "Status"],
        ["1", "Kredensial Valid, Akun Aktif, Perangkat Terpercaya", "Login Sukses", "Login Sukses", "Sesuai"],
        ["2", "Kredensial Valid, Akun Aktif, Perangkat Asing, Captcha Valid", "Login Sukses", "Login Sukses", "Sesuai"],
        ["3", "Kredensial Tidak Valid / Salah", "Login Ditolak", "Login Ditolak", "Sesuai"],
        ["4", "Akun Ditangguhkan (Suspended)", "Login Ditolak", "Login Ditolak", "Sesuai"],
        ["5", "Kredensial Valid, Perangkat Asing, Captcha Salah", "Login Ditolak", "Login Ditolak", "Sesuai"]
    ]

    test_table_shape = slide11.shapes.add_table(6, 5, Inches(0.7), Inches(1.9), Inches(11.5), Inches(4.0))
    test_table = test_table_shape.table

    test_table.columns[0].width = Inches(0.6)
    test_table.columns[1].width = Inches(5.2)
    test_table.columns[2].width = Inches(2.0)
    test_table.columns[3].width = Inches(2.0)
    test_table.columns[4].width = Inches(1.7)

    for r_idx, row in enumerate(test_data):
        for c_idx, val in enumerate(row):
            cell = test_table.cell(r_idx, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = 'Arial'
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER if c_idx != 1 else PP_ALIGN.LEFT
            if r_idx == 0:
                p.font.bold = True
                p.font.color.rgb = COLOR_HEADING
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xF8, 0xF5)
            else:
                p.font.color.rgb = COLOR_TEXT
                if c_idx in [2, 3]:
                    p.font.bold = True
                    if "Sukses" in val:
                        p.font.color.rgb = RGBColor(0x04, 0x78, 0x57)
                    else:
                        p.font.color.rgb = RGBColor(0xBE, 0x12, 0x3C)
                elif c_idx == 4:
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(0x04, 0x78, 0x57)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)

    # ==========================================================
    # SLIDE 12: Kesimpulan
    # ==========================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_background(slide12)

    # Centered card
    card12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.33), Inches(5.5))
    card12.fill.solid()
    card12.fill.fore_color.rgb = COLOR_CARD
    card12.line.color.rgb = COLOR_PRIMARY
    card12.line.width = Pt(2)

    accent12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.0), Inches(10.33), Inches(0.12))
    accent12.fill.solid()
    accent12.fill.fore_color.rgb = COLOR_PRIMARY
    accent12.line.fill.background()

    # Title
    t12_box = slide12.shapes.add_textbox(Inches(2.0), Inches(1.4), Inches(9.33), Inches(0.5))
    t12_p = t12_box.text_frame.paragraphs[0]
    t12_p.text = "KESIMPULAN AKHIR PROYEK"
    t12_p.font.name = 'Trebuchet MS'
    t12_p.font.size = Pt(18)
    t12_p.font.bold = True
    t12_p.font.color.rgb = COLOR_PRIMARY
    t12_p.alignment = PP_ALIGN.CENTER

    # Content
    c12_box = slide12.shapes.add_textbox(Inches(2.0), Inches(2.1), Inches(9.33), Inches(3.0))
    c12_tf = c12_box.text_frame
    c12_tf.word_wrap = True
    
    conclusions = [
        "1. Aplikasi MathSecure berhasil menjembatani teori matematika diskrit yang abstrak (Logika, Himpunan, Relasi, Kriptografi) dengan implementasi nyata pada sistem login berbasis React.",
        "2. Keberadaan visualisasi laboratorium interaktif (Diagram Venn SVG, Pohon Keputusan Dioda Menyala, Simulator Hashing) sangat membantu mempermudah pemahaman logika komputer mahasiswa.",
        "3. Penerapan sistem Context-Aware Authentication berhasil menyeimbangkan aspek keamanan siber terenkripsi kuat dengan kemudahan akses pengalaman login pengguna (UX)."
    ]
    for idx, c in enumerate(conclusions):
        cp = c12_tf.paragraphs[0] if idx == 0 else c12_tf.add_paragraph()
        cp.text = c
        cp.font.name = 'Arial'
        cp.font.size = Pt(12)
        cp.font.color.rgb = COLOR_TEXT
        cp.space_after = Pt(14)

    # Thanks
    thx_box = slide12.shapes.add_textbox(Inches(2.0), Inches(5.3), Inches(9.33), Inches(0.6))
    thx_p = thx_box.text_frame.paragraphs[0]
    thx_p.text = "Sekian & Terima Kasih - Ada Pertanyaan?"
    thx_p.font.name = 'Trebuchet MS'
    thx_p.font.size = Pt(18)
    thx_p.font.bold = True
    thx_p.font.color.rgb = COLOR_HEADING
    thx_p.alignment = PP_ALIGN.CENTER

    # Save presentation
    filename = "MathSecure_Presentasi_Matematika_Diskrit.pptx"
    prs.save(filename)
    print(f"PowerPoint presentation generated successfully: {filename}")

if __name__ == "__main__":
    build_presentation()
